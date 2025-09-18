import os
import re
import gc
import time
import json
import hashlib
import logging
import traceback
import tempfile
import shutil
import contextlib

from pathlib import Path
from typing import Callable, Dict, List, Set, Optional
from urllib.parse import urljoin, urlparse

import requests

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

# Service-friendly Canvas scraper module for running on a server (e.g., Heroku).
# Exposes: run_canvas_scrape_job(username, password, headless, status_callback)
#          -> {"input_path": str, "tmp_root": str}

START_URL = "https://canvas.cornell.edu"
COURSES_URL = "https://canvas.cornell.edu/courses"

ALLOWED_EXT_FOR_EXTRACTION = {
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".csv", ".txt", ".md", ".rtf"
}

MAX_LINKS_PER_COURSE = int(os.environ.get("MAX_LINKS_PER_COURSE", "250"))
MIN_TEXT_LEN_TO_RECORD = int(os.environ.get("MIN_TEXT_LEN_TO_RECORD", "80"))

# Hard caps to bound per-unit memory
MAX_PAGE_CHARS = int(os.environ.get("MAX_PAGE_CHARS", "50000"))   # per Canvas page write cap
MAX_FILE_CHARS = int(os.environ.get("MAX_FILE_CHARS", "200000"))  # per file write cap
# --- Warm-session settings (read by app.py edits) ----------------------------
REUSE_SESSION_ONLY = os.environ.get("OCEAN_REUSE_SESSION_ONLY", "0") == "1"
PERSIST_SESSION_DIR = os.environ.get("OCEAN_PERSIST_SESSION_DIR", "")
CHROME_PROFILE_DIR = os.environ.get("OCEAN_CHROME_PROFILE_DIR", "Default")
BRIDGE_COOKIES_PATH = os.environ.get("OCEAN_BRIDGE_COOKIES_PATH", "")



def trace_exc(msg="Exception"):
    logging.error("%s\n%s", msg, traceback.format_exc())


def make_dir(p: Path) -> Path:
    p.mkdir(parents=True, exist_ok=True)
    return p


def sanitize(name: str) -> str:
    return "".join([c if c.isalnum() or c in " _-" else "_" for c in (name or "").strip()]) or "Course"


def build_driver(headless: bool):
    chrome_opts = Options()
    if headless:
        chrome_opts.add_argument("--headless=new")
    # Memory/resource reduction flags
    chrome_opts.add_argument("--no-sandbox")
    chrome_opts.add_argument("--disable-gpu")
    chrome_opts.add_argument("--disable-dev-shm-usage")
    chrome_opts.add_argument("--window-size=1600,1200")
    chrome_opts.add_argument("--disable-background-networking")
    chrome_opts.add_argument("--disable-background-timer-throttling")
    chrome_opts.add_argument("--disable-renderer-backgrounding")
    chrome_opts.add_argument("--metrics-recording-only")
    chrome_opts.add_argument("--mute-audio")
    chrome_opts.add_argument("--no-first-run")
    chrome_opts.add_argument("--no-zygote")
    # Block heavy resources (images)
    chrome_opts.add_argument("--blink-settings=imagesEnabled=false")
    chrome_opts.add_experimental_option(
        "prefs",
        {
            "profile.managed_default_content_settings.images": 2,
            "download.prompt_for_download": False,
            "download.directory_upgrade": True,
            "plugins.always_open_pdf_externally": True,
        },
    )
    # Persisted Chrome profile so we keep the Canvas login warm between runs
    if PERSIST_SESSION_DIR:
        os.makedirs(PERSIST_SESSION_DIR, exist_ok=True)
        chrome_opts.add_argument(f"--user-data-dir={PERSIST_SESSION_DIR}")
        chrome_opts.add_argument(f"--profile-directory={CHROME_PROFILE_DIR}")
        # Helps avoid some profile-related flakiness in containers
        chrome_opts.add_argument("--disable-features=DialMediaRouteProvider")


    chrome_bin = (
        os.environ.get("GOOGLE_CHROME_BIN")
        or os.environ.get("CHROME_BIN")
        or os.environ.get("GOOGLE_CHROME_SHIM")
    )
    if chrome_bin:
        chrome_opts.binary_location = chrome_bin

    driver_path = (
        os.environ.get("CHROMEDRIVER_PATH")
        or os.environ.get("CHROMEWEBDRIVER")
    )
    try:
        if driver_path:
            if os.path.isdir(driver_path):
                driver_path = os.path.join(driver_path, "chromedriver")
            service = Service(executable_path=driver_path)
            return webdriver.Chrome(service=service, options=chrome_opts)
        return webdriver.Chrome(options=chrome_opts)
    except Exception:
        # Last resort
        return webdriver.Chrome(options=chrome_opts)


def _fallback_any_of(*conds):
    if hasattr(EC, "any_of"):
        return EC.any_of(*conds)

    class _AnyOf:
        def __init__(self, conditions):
            self.conditions = conditions

        def __call__(self, drv):
            for c in self.conditions:
                try:
                    res = c(drv)
                    if res:
                        return res
                except Exception:
                    pass
            return False

    return _AnyOf(conds)


def try_expand_all(driver, timeout: int = 5):
    try:
        wait = WebDriverWait(driver, timeout)
        btn = wait.until(EC.presence_of_element_located((By.ID, "expand_collapse_all")))
    except Exception:
        return False

    try:
        aria = (btn.get_attribute("aria-expanded") or "").strip().lower()
        de = (btn.get_attribute("data-expand") or "").strip().lower()
        should_click = (aria == "false") or (de == "false") or (not aria and not de)
        if should_click:
            btn.click()
            WebDriverWait(driver, timeout).until(
                lambda d: (
                    (btn.get_attribute("aria-expanded") or "").strip().lower() == "true"
                    or (btn.get_attribute("data-expand") or "").strip().lower() == "true"
                )
            )
        return True
    except Exception:
        return False


def scroll_to_bottom(driver, max_steps=20, pause=0.6):
    last_h = 0
    for _ in range(max_steps):
        driver.execute_script("window.scrollTo(0, document.body.scrollHeight);")
        time.sleep(pause)
        h = driver.execute_script("return document.body.scrollHeight;")
        if h == last_h:
            break
        last_h = h


def get_visible_text(driver) -> str:
    try:
        text = driver.execute_script(
            "return document.body && document.body.innerText ? document.body.innerText : ''"
        ) or ""
        if not text:
            text = driver.find_element(By.TAG_NAME, "body").text
        # normalize whitespace
        text = re.sub(r"[ \t\r\f\v]+", " ", text or "")
        text = re.sub(r"\n\s*\n+", "\n\n", text)
        return text.strip()
    except Exception:
        return ""


def get_course_title(driver) -> str:
    try:
        h1 = driver.find_element(By.XPATH, "//h1[contains(@class,'course-title') or contains(@class,'page-title')]")
        return sanitize(h1.text.strip())
    except Exception:
        try:
            return sanitize(driver.title.strip())
        except Exception:
            return "Course"


def normalize_link(href, course_id: str) -> str:
    try:
        raw = (href or "").strip()
    except Exception:
        return ""
    if not raw:
        return ""
    try:
        full = urljoin(START_URL + "/", raw)
    except Exception:
        full = raw
    try:
        p = urlparse(full)
    except Exception:
        return ""
    if not (p.scheme or "").startswith("http"):
        return ""

    host = (p.netloc or "").lower()
    if "canvas.cornell.edu" not in host:
        return ""

    p = p._replace(fragment="")
    url = p.geturl()
    path = (p.path or "")

    in_course = path.startswith(f"/courses/{course_id}") or f"/courses/{course_id}/" in path
    is_file = "/files/" in path
    if not (in_course or is_file):
        return ""

    bad = ["/login", "/conversations", "/calendar", "/profile", "/settings/profile", "/settings/notifications"]
    if any(b in path for b in bad):
        return ""

    return url
def _apply_initial_cookies(driver, cookies: List[Dict]) -> None:
    if not cookies:
        return
    try:
        driver.get(START_URL)
        for c in cookies:
            try:
                cookie = {
                    "name": c.get("name"),
                    "value": c.get("value"),
                    "path": c.get("path") or "/",
                    "domain": c.get("domain") or "canvas.cornell.edu",
                }
                if c.get("secure") is not None:
                    cookie["secure"] = bool(c.get("secure"))
                if c.get("expires") is not None:
                    cookie["expiry"] = int(c.get("expires"))
                driver.add_cookie(cookie)
            except Exception:
                pass
        # Reload with cookies in place
        driver.get(START_URL)
    except Exception:
        pass

def session_looks_logged_in(driver) -> bool:
    """
    Heuristic: load Canvas, see if we land on dashboard/courses and *not* on a login flow.
    """
    try:
        driver.get(START_URL)
        WebDriverWait(driver, 6).until(_fallback_any_of(
            EC.presence_of_element_located((By.ID, "dashboard")),
            EC.presence_of_element_located((By.XPATH, "//a[contains(@href, '/courses')]"))
        ))
    except Exception:
        # Even if wait fails, inspect URL to detect login pages
        pass

    try:
        cur = (driver.current_url or "").lower()
    except Exception:
        cur = ""

    # Explicit login patterns
    if "/login" in cur or "/saml" in cur:
        return False

    # Login form present?
    with contextlib.suppress(Exception):
        driver.find_element(By.XPATH, "//input[@name='j_username' or @id='username']")
        return False

    # Looks like we’re in a logged-in context (dashboard/courses/etc.)
    return True

def login_canvas(driver, username, password, status_callback: Callable[[str, str], None]):
    wait = WebDriverWait(driver, 30)
    driver.get(START_URL)

    try:
        btn = wait.until(EC.element_to_be_clickable((By.XPATH, "//a[contains(@href,'/login/saml') and contains(., 'Cornell')]")))
        btn.click()
    except Exception:
        driver.get(urljoin(START_URL, "/login/saml"))

    # Login form
    try:
        user_in = WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.XPATH, "//input[@name='j_username' or @id='username']")))
        pass_in = WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.XPATH, "//input[@name='j_password' or @id='password']")))
        user_in.clear(); user_in.send_keys(username)
        pass_in.clear(); pass_in.send_keys(password)
        try:
            login_btn = driver.find_element(By.XPATH, "//input[@type='submit' and (@name='_eventId_proceed' or @id='passwordbutton')]")
            login_btn.click()
        except Exception:
            pass_in.submit()
    except Exception:
        pass

    # Attempt to surface Duo code if present
    duo_code = None
    try:
        code_el = WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.CSS_SELECTOR, ".verification-code")))
        if code_el and code_el.text.strip():
            duo_code = code_el.text.strip()
    except Exception:
        pass

    if not duo_code:
        try:
            iframe = WebDriverWait(driver, 8).until(EC.presence_of_element_located((By.CSS_SELECTOR, "iframe#duo_iframe")))
            driver.switch_to.frame(iframe)
            try:
                code_el = WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.CSS_SELECTOR, ".verification-code")))
                if code_el and code_el.text.strip():
                    duo_code = code_el.text.strip()
            except Exception:
                pass
        except Exception:
            pass
        finally:
            with contextlib.suppress(Exception):
                driver.switch_to.default_content()

    if duo_code and status_callback:
        status_callback("duo", duo_code)
        status_callback("log", f"Duo pairing/verification code: {duo_code}")

    if status_callback:
        status_callback("status", "waiting_duo")

    # Allow time for Duo approval
    time.sleep(30)

    # Handle "shared device" prompts when present
    try:
        iframe = WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.CSS_SELECTOR, "iframe#duo_iframe")))
        driver.switch_to.frame(iframe)
        try:
            shared_button = WebDriverWait(driver, 10).until(
                EC.element_to_be_clickable((By.XPATH, "//button[contains(., 'No, other people use this device')]"))
            )
            shared_button.click()
        except Exception:
            pass
    except Exception:
        pass
    finally:
        with contextlib.suppress(Exception):
            driver.switch_to.default_content()

    with contextlib.suppress(Exception):
        el = driver.find_element(By.ID, "dont-trust-browser-button")
        el.click()

    with contextlib.suppress(Exception):
        handles = driver.window_handles
        if len(handles) > 1:
            driver.switch_to.window(handles[-1])

    robust_wait = WebDriverWait(driver, 60)
    dashboard_cond = EC.presence_of_element_located((By.ID, "dashboard"))
    courses_cond = EC.presence_of_element_located((By.XPATH, "//a[contains(@href, '/courses')]"))
    robust_wait.until(_fallback_any_of(dashboard_cond, courses_cond))

    if status_callback:
        status_callback("status", "logged_in")


def get_fall_2025_course_ids(driver) -> List[str]:
    ids: Set[str] = set()

    def _from_table(tbl) -> Set[str]:
        out: Set[str] = set()
        headers = tbl.find_elements(By.XPATH, ".//thead//th")
        term_idx = -1
        for i, th in enumerate(headers):
            txt = (th.text or "").strip().lower()
            if "term" in txt:
                term_idx = i

        rows = tbl.find_elements(By.XPATH, ".//tbody/tr")
        for row in rows:
            try:
                link = row.find_element(By.XPATH, ".//a[contains(@href, '/courses/') and not(contains(@href, '/users/'))]")
            except Exception:
                continue

            term_text = ""
            if term_idx >= 0:
                tds = row.find_elements(By.XPATH, ".//td")
                if term_idx < len(tds):
                    term_text = (tds[term_idx].text or "").strip()
            if not term_text:
                with contextlib.suppress(Exception):
                    cands = row.find_elements(By.XPATH, ".//td[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ','abcdefghijklmnopqrstuvwxyz'),'fall')]")
                    for td in cands:
                        txt = (td.text or "").strip()
                        if "fall" in txt.lower() and "2025" in txt:
                            term_text = txt
                            break
            if not term_text or "fall" not in term_text.lower() or "2025" not in term_text:
                continue

            href = link.get_attribute("href") or ""
            m = re.search(r"/courses/(\d+)", href)
            if m:
                out.add(m.group(1))
        return out

    driver.get(COURSES_URL)
    WebDriverWait(driver, 15).until(EC.presence_of_element_located((By.ID, "content")))
    tables = driver.find_elements(By.XPATH, "//table[.//thead//th[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'),'term')]]")
    for tbl in tables:
        ids |= _from_table(tbl)
    if ids:
        return list(ids)

    # Fallback: dashboard cards
    driver.get(START_URL)
    try:
        cards = WebDriverWait(driver, 10).until(
            EC.presence_of_all_elements_located((By.XPATH, "//a[contains(@href,'/courses/') and not(contains(@href,'/users/'))]"))
        )
    except Exception:
        cards = []

    cand_ids: Set[str] = set()
    for a in cards:
        with contextlib.suppress(Exception):
            href = a.get_attribute("href") or ""
            m = re.search(r"/courses/(\d+)", href)
            if m:
                cand_ids.add(m.group(1))

    def _verify_term(cid: str) -> bool:
        with contextlib.suppress(Exception):
            driver.get(f"{START_URL}/courses/{cid}/settings")
            term_el = WebDriverWait(driver, 8).until(
                EC.presence_of_element_located((
                    By.XPATH,
                    "//*[self::label or self::div or self::span][contains(., 'Term')]/following::span[1] | //*[contains(., 'Term')]/following::*[1]"
                ))
            )
            term_txt = (term_el.text or "").strip().lower()
            return ("fall" in term_txt) and ("2025" in term_txt)
        return False

    for cid in cand_ids:
        if _verify_term(cid):
            ids.add(cid)

    return list(ids)


def append_header_and_streamed_file(input_txt_path: Path, course_name: str, src_path: Path, origin_url: str):
    header = f"--- Scraped from {course_name} at {origin_url} ---\n"
    with open(input_txt_path, "a", encoding="utf-8", errors="ignore") as out_f:
        out_f.write(header)
        with open(src_path, "r", encoding="utf-8", errors="ignore") as in_f:
            shutil.copyfileobj(in_f, out_f, length=64 * 1024)  # copy in chunks
        out_f.write("\n\n")


def collect_in_course_links(driver, course_id: str) -> Set[str]:
    # Avoid many WebElement objects; pull hrefs via JS
    hrefs: List[str] = []
    with contextlib.suppress(Exception):
        hrefs = driver.execute_script(
            "return Array.from(document.querySelectorAll('a[href]')).map(a => a.href);"
        ) or []
    urls: Set[str] = set()
    for raw in hrefs:
        norm = normalize_link(raw, course_id)
        if norm:
            urls.add(norm)
    return urls


def download_file_from_canvas(driver, file_url: str, download_dir: Path, session: requests.Session) -> Optional[Path]:
    driver.get(file_url)
    try_expand_all(driver, timeout=3)
    time.sleep(1.2)

    download_link = None
    try:
        download_link = WebDriverWait(driver, 6).until(
            EC.element_to_be_clickable((By.XPATH, "//a[contains(text(), 'Download') or contains(@href, '/download')]"))
        )
    except Exception:
        with contextlib.suppress(Exception):
            download_link = driver.find_element(By.XPATH, "//a[contains(@class, 'btn') and contains(@href, '/download')]")
    if not download_link:
        return None

    download_url = download_link.get_attribute("href") or file_url

    try:
        title_el = driver.find_element(By.XPATH, "//h1 | //h2")
        filename = (title_el.text or "").strip()
    except Exception:
        filename = ""

    if not filename or not any(filename.lower().endswith(ext) for ext in ALLOWED_EXT_FOR_EXTRACTION):
        path_name = os.path.basename(urlparse(download_url).path)
        if path_name and "." in path_name:
            filename = path_name
        else:
            filename = f"file_{int(time.time())}.pdf"

    filename = sanitize(filename)
    suf = os.path.splitext(filename)[1].lower()
    if not any(suf.endswith(ext) for ext in ALLOWED_EXT_FOR_EXTRACTION):
        filename += ".pdf"

    file_path = download_dir / filename

    # Carry cookies over for authorized download
    with contextlib.suppress(Exception):
        for cookie in driver.get_cookies():
            with contextlib.suppress(Exception):
                session.cookies.set(cookie["name"], cookie["value"], domain=cookie.get("domain"))

    resp = session.get(download_url, stream=True, timeout=60)
    resp.raise_for_status()
    with open(file_path, "wb") as f:
        for chunk in resp.iter_content(8192):
            if not chunk:
                continue
            f.write(chunk)

    return file_path


def _stream_pdf_to_file(pdf_path: Path, tmp_text_path: Path) -> int:
    written = 0
    # Prefer PyMuPDF (low memory) if available
    try:
        import fitz  # PyMuPDF
        with fitz.open(str(pdf_path)) as doc, open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for page in doc:
                txt = page.get_text("text") or ""
                if not txt:
                    continue
                if MAX_FILE_CHARS and written + len(txt) > MAX_FILE_CHARS:
                    txt = txt[: max(0, MAX_FILE_CHARS - written)]
                out.write(txt + "\n")
                written += len(txt)
                if written >= MAX_FILE_CHARS:
                    break
        return written
    except Exception:
        pass

    # Fallback to PyPDF2, still stream page-by-page
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))
        with open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for page in reader.pages:
                try:
                    txt = (page.extract_text() or "")
                except Exception:
                    txt = ""
                if not txt:
                    continue
                if MAX_FILE_CHARS and written + len(txt) > MAX_FILE_CHARS:
                    txt = txt[: max(0, MAX_FILE_CHARS - written)]
                out.write(txt + "\n")
                written += len(txt)
                if written >= MAX_FILE_CHARS:
                    break
        return written
    except Exception:
        return 0


def _stream_docx_to_file(docx_path: Path, tmp_text_path: Path) -> int:
    written = 0
    try:
        from docx import Document as DocxDocument  # python-docx
        doc = DocxDocument(str(docx_path))
        with open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for p in doc.paragraphs:
                txt = (p.text or "")
                if not txt:
                    continue
                if MAX_FILE_CHARS and written + len(txt) > MAX_FILE_CHARS:
                    txt = txt[: max(0, MAX_FILE_CHARS - written)]
                out.write(txt + "\n")
                written += len(txt)
                if written >= MAX_FILE_CHARS:
                    break
        return written
    except Exception:
        # Fallback avoided to keep memory low
        return 0


def _stream_pptx_to_file(pptx_path: Path, tmp_text_path: Path) -> int:
    written = 0
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        with open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        txt = shape.text
                        if MAX_FILE_CHARS and written + len(txt) > MAX_FILE_CHARS:
                            txt = txt[: max(0, MAX_FILE_CHARS - written)]
                        out.write(txt + "\n")
                        written += len(txt)
                        if written >= MAX_FILE_CHARS:
                            return written
        return written
    except Exception:
        return 0


def _stream_xlsx_to_file(xlsx_path: Path, tmp_text_path: Path) -> int:
    # Use openpyxl read-only mode: constant memory
    written = 0
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path), read_only=True, data_only=True)
        with open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = ",".join("" if v is None else str(v) for v in row)
                    if MAX_FILE_CHARS and written + len(line) > MAX_FILE_CHARS:
                        line = line[: max(0, MAX_FILE_CHARS - written)]
                    out.write(line + "\n")
                    written += len(line)
                    if written >= MAX_FILE_CHARS:
                        wb.close()
                        return written
        wb.close()
        return written
    except Exception:
        return 0


def _stream_txt_like_to_file(path: Path, tmp_text_path: Path) -> int:
    written = 0
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as inp, open(tmp_text_path, "w", encoding="utf-8", errors="ignore") as out:
            for line in inp:
                if MAX_FILE_CHARS and written + len(line) > MAX_FILE_CHARS:
                    line = line[: max(0, MAX_FILE_CHARS - written)]
                out.write(line)
                written += len(line)
                if written >= MAX_FILE_CHARS:
                    break
        return written
    except Exception:
        return 0


def stream_extract_file_to_temp(path: Path, tmp_text_path: Path) -> int:
    suf = (path.suffix or "").lower()
    if suf == ".pdf":
        return _stream_pdf_to_file(path, tmp_text_path)
    if suf == ".docx":
        return _stream_docx_to_file(path, tmp_text_path)
    if suf == ".pptx":
        return _stream_pptx_to_file(path, tmp_text_path)
    if suf in {".txt", ".md", ".csv"}:
        return _stream_txt_like_to_file(path, tmp_text_path)
    if suf in {".xlsx"}:
        return _stream_xlsx_to_file(path, tmp_text_path)
    if suf in {".doc"}:
        # Not supported in streaming mode; skip to avoid memory spikes
        return 0
    if suf in {".xls"}:
        # Skip to avoid loading entire workbook; could be supported with xlrd if needed
        return 0
    # Unknown formats: skip
    return 0


def _write_tmp_text(text: str) -> Path:
    # Utility: write small (capped) page text to a temp file for streaming append
    tmpf = tempfile.NamedTemporaryFile("w+", delete=False, encoding="utf-8")
    path = Path(tmpf.name)
    try:
        tmpf.write(text)
        tmpf.write("\n")
        tmpf.flush()
        tmpf.close()
        return path
    except Exception:
        with contextlib.suppress(Exception):
            tmpf.close()
        return path


def _push_snippet(status_callback: Callable[[str, str], None], kind: str, where: str, text: str, course: str):
    if not status_callback or not text:
        return
    # normalize whitespace and cap preview length
    snippet = re.sub(r"\s+", " ", text).strip()
    if len(snippet) > 600:
        snippet = snippet[:600] + "..."
    status_callback("snippet", f"[{course}] {kind} {where}: {snippet}")


def crawl_course(driver, course_id: str, input_txt_path: Path, status_callback: Callable[[str, str], None]):
    base = f"{START_URL}/courses/{course_id}"
    seeds = [
        base,
        f"{base}/assignments",
        f"{base}/modules",
        f"{base}/assignments/syllabus",
        f"{base}/grades",
        f"{base}/announcements",
    ]

    driver.get(base)
    try_expand_all(driver, 5)
    time.sleep(1.0)

    course_name = get_course_title(driver)
    course_dir = input_txt_path.parent / sanitize(course_name)
    make_dir(course_dir)

    session = requests.Session()

    # First: modules quick harvest of files
    with contextlib.suppress(Exception):
        driver.get(f"{base}/modules")
        try_expand_all(driver, 5)
        time.sleep(1.0)
        links = driver.find_elements(By.XPATH, "//a[contains(@href, '/files/') or contains(@href, '.pdf') or contains(@href, '.docx') or contains(@href, '.pptx') or contains(@href, '.xlsx') or contains(@href, '.csv')]")
        seen_files: Set[str] = set()
        for a in links:
            href = a.get_attribute("href") or ""
            if "/files/" in href and href not in seen_files:
                seen_files.add(href)
                with contextlib.suppress(Exception):
                    p = download_file_from_canvas(driver, href, course_dir, session)
                    if p and p.exists():
                        with tempfile.NamedTemporaryFile("w+", delete=False, encoding="utf-8") as tmpf:
                            tmp_txt_path = Path(tmpf.name)
                        try:
                            written = stream_extract_file_to_temp(p, tmp_txt_path)
                            if written >= MIN_TEXT_LEN_TO_RECORD:
                                append_header_and_streamed_file(input_txt_path, course_name, tmp_txt_path, f"FILE: {p.name} ({href})")
                                if status_callback:
                                    status_callback("log", f"[file] saved {p.name} ({written} chars) from {href}")
                                # Push live snippet
                                try:
                                    with open(tmp_txt_path, "r", encoding="utf-8", errors="ignore") as r:
                                        content = r.read(MAX_PAGE_CHARS)
                                    _push_snippet(status_callback, "FILE", p.name, content, course_name)
                                except Exception:
                                    pass
                        finally:
                            with contextlib.suppress(Exception):
                                tmp_txt_path.unlink()

    visited_pages_h: Set[bytes] = set()
    visited_files_h: Set[bytes] = set()
    queue: List[str] = list(seeds)
    steps = 0

    while queue and len(visited_pages_h) < MAX_LINKS_PER_COURSE:
        url = queue.pop(0)
        h = hashlib.md5(url.encode("utf-8")).digest()
        if h in visited_pages_h:
            continue
        visited_pages_h.add(h)

        try:
            driver.get(url)
            try_expand_all(driver, 5)
            time.sleep(0.6)
            scroll_to_bottom(driver, 12, 0.3)

            # Page text with cap
            page_text = get_visible_text(driver)
            if page_text:
                truncated = False
                if MAX_PAGE_CHARS and len(page_text) > MAX_PAGE_CHARS:
                    page_text = page_text[:MAX_PAGE_CHARS]
                    truncated = True

                if len(page_text) >= MIN_TEXT_LEN_TO_RECORD:
                    tmp_txt_path = _write_tmp_text(page_text)
                    try:
                        append_header_and_streamed_file(input_txt_path, course_name, tmp_txt_path, url)
                    finally:
                        with contextlib.suppress(Exception):
                            Path(tmp_txt_path).unlink()

                    if status_callback:
                        status_callback("log", f"[page] {url} -> {len(page_text)} chars{' (truncated)' if truncated else ''}")
                    # Live snippet from the page
                    _push_snippet(status_callback, "PAGE", url, page_text, course_name)

            # Discover links
            links = collect_in_course_links(driver, course_id)
            for link in links:
                is_file = (
                    ("/files/" in link)
                    or any(link.lower().endswith(ext) for ext in ALLOWED_EXT_FOR_EXTRACTION)
                    or "/download" in link.lower()
                )
                if is_file:
                    hf = hashlib.md5(link.encode("utf-8")).digest()
                    if hf in visited_files_h:
                        continue
                    visited_files_h.add(hf)
                    with contextlib.suppress(Exception):
                        p = download_file_from_canvas(driver, link, course_dir, session)
                        if p and p.exists():
                            with tempfile.NamedTemporaryFile("w+", delete=False, encoding="utf-8") as tmpf:
                                tmp_txt_path = Path(tmpf.name)
                            try:
                                written = stream_extract_file_to_temp(p, tmp_txt_path)
                                if written >= MIN_TEXT_LEN_TO_RECORD:
                                    append_header_and_streamed_file(input_txt_path, course_name, tmp_txt_path, f"FILE: {p.name} ({link})")
                                    if status_callback:
                                        status_callback("log", f"[file] saved {p.name} ({written} chars) from {link}")
                                    # Snippet for file content
                                    try:
                                        with open(tmp_txt_path, "r", encoding="utf-8", errors="ignore") as r:
                                            content = r.read(MAX_PAGE_CHARS)
                                        _push_snippet(status_callback, "FILE", p.name, content, course_name)
                                    except Exception:
                                        pass
                            finally:
                                with contextlib.suppress(Exception):
                                    tmp_txt_path.unlink()
                else:
                    if len(visited_pages_h) + len(queue) < MAX_LINKS_PER_COURSE:
                        queue.append(link)

        except Exception:
            pass

        steps += 1
        if steps % 10 == 0:
            gc.collect()
            if status_callback:
                status_callback("log", f"Crawled {len(visited_pages_h)} pages and {len(visited_files_h)} file endpoints in course {course_id}")

    if status_callback:
        status_callback("log", f"[course done] {course_id}: {len(visited_pages_h)} pages, {len(visited_files_h)} file endpoints")


def run_course_crawl(driver, course_id: str, input_txt_path: Path, status_callback: Callable[[str, str], None]):
    base = f"{START_URL}/courses/{course_id}"
    driver.get(base)
    try_expand_all(driver, timeout=5)
    time.sleep(1.0)

    course_name = get_course_title(driver)
    course_dir = input_txt_path.parent / sanitize(course_name)
    make_dir(course_dir)

    # Quick file harvest then full crawl via BFS with snippets/logs
    session = requests.Session()
    with contextlib.suppress(Exception):
        driver.get(f"{base}/modules")
        try_expand_all(driver, 5)
        time.sleep(1.0)
        links = driver.find_elements(By.XPATH, "//a[contains(@href, '/files/') or contains(@href, '.pdf') or contains(@href, '.docx') or contains(@href, '.pptx') or contains(@href, '.xlsx') or contains(@href, '.csv')]")
        for a in links:
            href = a.get_attribute("href") or ""
            if "/files/" in href:
                with contextlib.suppress(Exception):
                    p = download_file_from_canvas(driver, href, course_dir, session)
                    if p and p.exists():
                        with tempfile.NamedTemporaryFile("w+", delete=False, encoding="utf-8") as tmpf:
                            tmp_txt_path = Path(tmpf.name)
                        try:
                            written = stream_extract_file_to_temp(p, tmp_txt_path)
                            if written >= MIN_TEXT_LEN_TO_RECORD:
                                append_header_and_streamed_file(input_txt_path, course_name, tmp_txt_path, f"FILE: {p.name} ({href})")
                                if status_callback:
                                    status_callback("log", f"[file-prefetch] saved {p.name} ({written} chars) from {href}")
                                try:
                                    with open(tmp_txt_path, "r", encoding="utf-8", errors="ignore") as r:
                                        content = r.read(MAX_PAGE_CHARS)
                                    _push_snippet(status_callback, "FILE", p.name, content, course_name)
                                except Exception:
                                    pass
                        finally:
                            with contextlib.suppress(Exception):
                                tmp_txt_path.unlink()

    # BFS across main tabs and discovered links with live updates
    crawl_course(driver, course_id, input_txt_path, status_callback)


def run_canvas_scrape_job(username: str, password: str, headless: bool, status_callback: Callable[[str, str], None], initial_cookies: Optional[List[Dict]] = None) -> Dict:

    Run a complete Canvas scrape job and return the aggregated input file path and temp root for later cleanup.
    """
    tmp_root = Path(tempfile.mkdtemp(prefix="canvas_job_"))
    input_txt = tmp_root / "input.txt"
    input_txt.write_text("Canvas Raw Input (aggregated page and file text)\n", encoding="utf-8")

    driver = build_driver(headless=headless)
    try:
        # First, try to reuse an existing warm session
        if status_callback:
            status_callback("status", "checking_session")

        logged_in = session_looks_logged_in(driver)
            # If we were given cookies from the bridge (user logged in), apply them now
        if not logged_in:
            if initial_cookies and isinstance(initial_cookies, list):
                _apply_initial_cookies(driver, initial_cookies)
                logged_in = session_looks_logged_in(driver)
            elif BRIDGE_COOKIES_PATH and os.path.exists(BRIDGE_COOKIES_PATH):
                try:
                    with open(BRIDGE_COOKIES_PATH, "r", encoding="utf-8") as f:
                        lst = json.load(f)
                    if isinstance(lst, list):
                        _apply_initial_cookies(driver, lst)
                        logged_in = session_looks_logged_in(driver)
                except Exception:
                    pass

        if not logged_in:
            if REUSE_SESSION_ONLY or not (username and password):
                # In reuse-only mode or no creds provided: do NOT log in
                if status_callback:
                    status_callback("log", "No valid Canvas session and reuse-only mode set -> skipping login")
                    status_callback("status", "login_required")
                return {"input_path": "", "tmp_root": str(tmp_root)}

            # Allowed to log in (manual run)
            if status_callback:
                status_callback("status", "logging_in")
            login_canvas(driver, username, password, status_callback)
        else:
            if status_callback:
                status_callback("status", "logged_in")

        if status_callback:
            status_callback("status", "discovering_courses")

        course_ids = get_fall_2025_course_ids(driver)

        if status_callback:
            status_callback("log", f"Found {len(course_ids)} Fall 2025 courses: {course_ids}")

        for cid in course_ids:
            if status_callback:
                status_callback("log", f"Processing course {cid}")
            run_course_crawl(driver, cid, input_txt, status_callback)

        if status_callback:
            status_callback("status", "completed")

        return {"input_path": str(input_txt), "tmp_root": str(tmp_root)}
    except Exception as e:
        if status_callback:
            status_callback("log", f"error: {e}")
            status_callback("status", "failed")
        return {"input_path": "", "tmp_root": str(tmp_root)}
    finally:
        with contextlib.suppress(Exception):
            driver.quit()

    # Do NOT delete tmp_root here; app.py will clean up after embedding



