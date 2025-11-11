#!/usr/bin/env python3
"""
Comprehensive Canvas Course Scraper
====================================

This scraper is designed to extract ALL content from a Canvas course, including:
- Syllabus pages and course information
- Module content (PDFs, PowerPoints, Word docs, etc.)
- Files section
- Assignments
- Announcements
- Canvas pages
- External links (noted but not followed outside Canvas domain)

The scraper is adaptable to different course structures and handles:
- Courses with extensive module organization (like AEP 4200, GOVT 1111)
- Courses with minimal modules but extensive Files section (like AEP 4230)
- Courses with external integrations like Google Drive (like MATH 4220)
- Various file types: PDF, PPTX, DOCX, ZIP, etc.

Usage:
    python canvas_course_scraper.py <course_url>
    
    Example: python canvas_course_scraper.py https://canvas.cornell.edu/courses/80403
"""

import os
import sys
import time
import json
import re
from urllib.parse import urljoin, urlparse, parse_qs
from pathlib import Path

from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException, NoSuchElementException
from bs4 import BeautifulSoup
import requests

# For extracting text from different file types
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    print("Warning: PDF libraries not installed. Install with: pip install PyPDF2 pdfplumber")

try:
    from pptx import Presentation
except ImportError:
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")

try:
    from docx import Document
except ImportError:
    print("Warning: python-docx not installed. Install with: pip install python-docx")


class CanvasCourseScr aper:
    """Comprehensive Canvas course scraper that adapts to different course structures."""
    
    def __init__(self, course_url, output_dir="scraped_course"):
        """
        Initialize the scraper.
        
        Args:
            course_url: Full URL to the Canvas course (e.g., https://canvas.cornell.edu/courses/80403)
            output_dir: Directory to save scraped content
        """
        self.course_url = course_url.rstrip('/')
        self.course_id = self._extract_course_id(course_url)
        self.base_url = self._get_base_url(course_url)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Track visited URLs to avoid duplicates
        self.visited_urls = set()
        self.visited_files = set()
        
        # Store scraped content metadata
        self.content_index = {
            'course_url': course_url,
            'course_id': self.course_id,
            'sections': {},
            'files': [],
            'external_links': []
        }
        
        # Initialize browser (assumes already logged in)
        self.driver = None
        
    def _extract_course_id(self, url):
        """Extract course ID from URL."""
        match = re.search(r'/courses/(\d+)', url)
        if match:
            return match.group(1)
        raise ValueError(f"Could not extract course ID from URL: {url}")
    
    def _get_base_url(self, url):
        """Get base URL (scheme + netloc)."""
        parsed = urlparse(url)
        return f"{parsed.scheme}://{parsed.netloc}"
    
    def _is_canvas_url(self, url):
        """Check if URL is within Canvas domain."""
        if not url:
            return False
        parsed = urlparse(url)
        return 'canvas' in parsed.netloc or parsed.netloc == ''
    
    def _normalize_filename(self, filename):
        """Normalize filename to be filesystem-safe."""
        # Remove or replace invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Limit length
        if len(filename) > 200:
            name, ext = os.path.splitext(filename)
            filename = name[:200-len(ext)] + ext
        return filename
    
    def init_browser(self):
        """Initialize Selenium browser (assumes cookies/session already set)."""
        options = webdriver.ChromeOptions()
        options.add_argument('--headless')
        options.add_argument('--no-sandbox')
        options.add_argument('--disable-dev-shm-usage')
        self.driver = webdriver.Chrome(options=options)
        self.driver.implicitly_wait(10)
    
    def close_browser(self):
        """Close the browser."""
        if self.driver:
            self.driver.quit()
    
    def scrape_all(self):
        """Main method to scrape all course content."""
        print(f"Starting comprehensive scrape of course {self.course_id}")
        print(f"Output directory: {self.output_dir}")
        
        try:
            self.init_browser()
            
            # Navigate to course home
            self.driver.get(self.course_url)
            time.sleep(2)
            
            # Scrape different sections
            print("\n[1/7] Scraping course home page...")
            self.scrape_home_page()
            
            print("\n[2/7] Scraping syllabus...")
            self.scrape_syllabus()
            
            print("\n[3/7] Scraping modules...")
            self.scrape_modules()
            
            print("\n[4/7] Scraping files section...")
            self.scrape_files_section()
            
            print("\n[5/7] Scraping assignments...")
            self.scrape_assignments()
            
            print("\n[6/7] Scraping announcements...")
            self.scrape_announcements()
            
            print("\n[7/7] Scraping additional pages...")
            self.scrape_pages()
            
            # Save content index
            self._save_content_index()
            
            print(f"\n✓ Scraping complete! Content saved to: {self.output_dir}")
            print(f"  - Total files downloaded: {len(self.visited_files)}")
            print(f"  - Total pages scraped: {len(self.visited_urls)}")
            print(f"  - External links noted: {len(self.content_index['external_links'])}")
            
        finally:
            self.close_browser()
    
    def scrape_home_page(self):
        """Scrape the course home page."""
        try:
            self.driver.get(self.course_url)
            time.sleep(2)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            # Extract course title
            title_elem = soup.find('h1') or soup.find('title')
            course_title = title_elem.get_text(strip=True) if title_elem else "Course"
            
            # Save home page content
            home_dir = self.output_dir / "home"
            home_dir.mkdir(exist_ok=True)
            
            with open(home_dir / "home_page.html", 'w', encoding='utf-8') as f:
                f.write(self.driver.page_source)
            
            # Extract text content
            content = soup.get_text(separator='\n', strip=True)
            with open(home_dir / "home_page.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            
            self.content_index['course_title'] = course_title
            self.content_index['sections']['home'] = {
                'files': ['home/home_page.html', 'home/home_page.txt']
            }
            
            print(f"  ✓ Saved home page: {course_title}")
            
        except Exception as e:
            print(f"  ✗ Error scraping home page: {e}")
    
    def scrape_syllabus(self):
        """Scrape the syllabus page."""
        try:
            syllabus_url = f"{self.course_url}/assignments/syllabus"
            self.driver.get(syllabus_url)
            time.sleep(2)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            # Save syllabus
            syllabus_dir = self.output_dir / "syllabus"
            syllabus_dir.mkdir(exist_ok=True)
            
            with open(syllabus_dir / "syllabus.html", 'w', encoding='utf-8') as f:
                f.write(self.driver.page_source)
            
            # Extract text
            content = soup.get_text(separator='\n', strip=True)
            with open(syllabus_dir / "syllabus.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Look for syllabus links/files
            links = soup.find_all('a', href=True)
            for link in links:
                href = link['href']
                if 'syllabus' in link.get_text().lower() or 'syllabus' in href.lower():
                    if href.endswith(('.pdf', '.docx', '.doc')):
                        self._download_file(href, syllabus_dir)
            
            self.content_index['sections']['syllabus'] = {
                'files': ['syllabus/syllabus.html', 'syllabus/syllabus.txt']
            }
            
            print(f"  ✓ Saved syllabus")
            
        except Exception as e:
            print(f"  ✗ Error scraping syllabus: {e}")
    
    def scrape_modules(self):
        """Scrape all modules and their content."""
        try:
            modules_url = f"{self.course_url}/modules"
            self.driver.get(modules_url)
            time.sleep(2)
            
            # Click "Expand All" button if it exists
            try:
                expand_btn = self.driver.find_element(By.XPATH, "//button[contains(text(), 'Expand All')]")
                expand_btn.click()
                time.sleep(2)
            except NoSuchElementException:
                print("  No 'Expand All' button found, modules may already be expanded")
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            # Find all modules
            modules = soup.find_all('div', class_=lambda x: x and 'context_module' in x)
            
            if not modules:
                # Try alternative structure
                modules = soup.find_all('div', attrs={'role': 'region', 'aria-label': lambda x: x and 'module' in x.lower()})
            
            print(f"  Found {len(modules)} modules")
            
            modules_dir = self.output_dir / "modules"
            modules_dir.mkdir(exist_ok=True)
            
            module_data = []
            
            for idx, module in enumerate(modules, 1):
                try:
                    # Get module name
                    module_name = module.find(['h2', 'h3', 'span'], class_=lambda x: x and ('header' in x or 'title' in x))
                    if module_name:
                        module_name = module_name.get_text(strip=True)
                    else:
                        module_name = f"Module_{idx}"
                    
                    module_name = self._normalize_filename(module_name)
                    print(f"  Processing module: {module_name}")
                    
                    module_dir = modules_dir / module_name
                    module_dir.mkdir(exist_ok=True)
                    
                    # Find all items in this module
                    items = module.find_all('a', href=True)
                    
                    module_items = []
                    for item in items:
                        href = item['href']
                        item_text = item.get_text(strip=True)
                        
                        if not href or href.startswith('#'):
                            continue
                        
                        # Make absolute URL
                        if not href.startswith('http'):
                            href = urljoin(self.base_url, href)
                        
                        # Skip if already visited
                        if href in self.visited_urls:
                            continue
                        
                        print(f"    - {item_text}")
                        
                        # Determine item type and handle accordingly
                        if '/files/' in href:
                            # File item
                            file_path = self._download_file(href, module_dir)
                            if file_path:
                                module_items.append({
                                    'type': 'file',
                                    'name': item_text,
                                    'path': str(file_path.relative_to(self.output_dir))
                                })
                        elif '/pages/' in href:
                            # Canvas page
                            page_path = self._scrape_page(href, module_dir, item_text)
                            if page_path:
                                module_items.append({
                                    'type': 'page',
                                    'name': item_text,
                                    'path': str(page_path.relative_to(self.output_dir))
                                })
                        elif '/assignments/' in href:
                            # Assignment
                            assign_path = self._scrape_assignment(href, module_dir, item_text)
                            if assign_path:
                                module_items.append({
                                    'type': 'assignment',
                                    'name': item_text,
                                    'path': str(assign_path.relative_to(self.output_dir))
                                })
                        elif self._is_canvas_url(href):
                            # Other Canvas content
                            self.visited_urls.add(href)
                            module_items.append({
                                'type': 'link',
                                'name': item_text,
                                'url': href
                            })
                        else:
                            # External link
                            self.content_index['external_links'].append({
                                'name': item_text,
                                'url': href,
                                'context': f'Module: {module_name}'
                            })
                            module_items.append({
                                'type': 'external_link',
                                'name': item_text,
                                'url': href
                            })
                    
                    module_data.append({
                        'name': module_name,
                        'items': module_items
                    })
                    
                except Exception as e:
                    print(f"  ✗ Error processing module {idx}: {e}")
                    continue
            
            # Save module index
            with open(modules_dir / "modules_index.json", 'w', encoding='utf-8') as f:
                json.dump(module_data, f, indent=2)
            
            self.content_index['sections']['modules'] = module_data
            print(f"  ✓ Processed {len(module_data)} modules")
            
        except Exception as e:
            print(f"  ✗ Error scraping modules: {e}")
    
    def scrape_files_section(self):
        """Scrape the Files section (important for courses like AEP 4230)."""
        try:
            files_url = f"{self.course_url}/files"
            self.driver.get(files_url)
            time.sleep(3)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            files_dir = self.output_dir / "files_section"
            files_dir.mkdir(exist_ok=True)
            
            # Find all file links
            file_links = soup.find_all('a', href=lambda x: x and '/files/' in x)
            
            print(f"  Found {len(file_links)} files in Files section")
            
            downloaded_files = []
            for link in file_links:
                try:
                    href = link['href']
                    file_name = link.get_text(strip=True)
                    
                    if not href.startswith('http'):
                        href = urljoin(self.base_url, href)
                    
                    # Skip if already downloaded
                    if href in self.visited_files:
                        continue
                    
                    print(f"    - {file_name}")
                    file_path = self._download_file(href, files_dir)
                    
                    if file_path:
                        downloaded_files.append({
                            'name': file_name,
                            'path': str(file_path.relative_to(self.output_dir))
                        })
                        
                except Exception as e:
                    print(f"    ✗ Error downloading file: {e}")
                    continue
            
            self.content_index['sections']['files'] = downloaded_files
            print(f"  ✓ Downloaded {len(downloaded_files)} files from Files section")
            
        except Exception as e:
            print(f"  ✗ Error scraping files section: {e}")
    
    def scrape_assignments(self):
        """Scrape assignments."""
        try:
            assignments_url = f"{self.course_url}/assignments"
            self.driver.get(assignments_url)
            time.sleep(2)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            assignments_dir = self.output_dir / "assignments"
            assignments_dir.mkdir(exist_ok=True)
            
            # Find assignment links
            assignment_links = soup.find_all('a', href=lambda x: x and '/assignments/' in x)
            
            print(f"  Found {len(assignment_links)} assignments")
            
            assignments_data = []
            for link in assignment_links:
                try:
                    href = link['href']
                    assign_name = link.get_text(strip=True)
                    
                    if not href.startswith('http'):
                        href = urljoin(self.base_url, href)
                    
                    if href in self.visited_urls:
                        continue
                    
                    print(f"    - {assign_name}")
                    assign_path = self._scrape_assignment(href, assignments_dir, assign_name)
                    
                    if assign_path:
                        assignments_data.append({
                            'name': assign_name,
                            'path': str(assign_path.relative_to(self.output_dir))
                        })
                        
                except Exception as e:
                    print(f"    ✗ Error scraping assignment: {e}")
                    continue
            
            self.content_index['sections']['assignments'] = assignments_data
            print(f"  ✓ Scraped {len(assignments_data)} assignments")
            
        except Exception as e:
            print(f"  ✗ Error scraping assignments: {e}")
    
    def scrape_announcements(self):
        """Scrape course announcements."""
        try:
            announcements_url = f"{self.course_url}/announcements"
            self.driver.get(announcements_url)
            time.sleep(2)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            announcements_dir = self.output_dir / "announcements"
            announcements_dir.mkdir(exist_ok=True)
            
            # Save announcements page
            with open(announcements_dir / "announcements.html", 'w', encoding='utf-8') as f:
                f.write(self.driver.page_source)
            
            content = soup.get_text(separator='\n', strip=True)
            with open(announcements_dir / "announcements.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            
            self.content_index['sections']['announcements'] = {
                'files': ['announcements/announcements.html', 'announcements/announcements.txt']
            }
            
            print(f"  ✓ Saved announcements")
            
        except Exception as e:
            print(f"  ✗ Error scraping announcements: {e}")
    
    def scrape_pages(self):
        """Scrape additional Canvas pages."""
        try:
            pages_url = f"{self.course_url}/pages"
            self.driver.get(pages_url)
            time.sleep(2)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            pages_dir = self.output_dir / "pages"
            pages_dir.mkdir(exist_ok=True)
            
            # Find page links
            page_links = soup.find_all('a', href=lambda x: x and '/pages/' in x)
            
            print(f"  Found {len(page_links)} pages")
            
            pages_data = []
            for link in page_links:
                try:
                    href = link['href']
                    page_name = link.get_text(strip=True)
                    
                    if not href.startswith('http'):
                        href = urljoin(self.base_url, href)
                    
                    if href in self.visited_urls:
                        continue
                    
                    print(f"    - {page_name}")
                    page_path = self._scrape_page(href, pages_dir, page_name)
                    
                    if page_path:
                        pages_data.append({
                            'name': page_name,
                            'path': str(page_path.relative_to(self.output_dir))
                        })
                        
                except Exception as e:
                    print(f"    ✗ Error scraping page: {e}")
                    continue
            
            self.content_index['sections']['pages'] = pages_data
            print(f"  ✓ Scraped {len(pages_data)} pages")
            
        except Exception as e:
            print(f"  ✗ Error scraping pages: {e}")
    
    def _download_file(self, url, save_dir):
        """
        Download a file from Canvas.
        
        Args:
            url: URL to the file
            save_dir: Directory to save the file
            
        Returns:
            Path to downloaded file, or None if failed
        """
        try:
            if url in self.visited_files:
                return None
            
            self.visited_files.add(url)
            
            # Navigate to file page to get download link
            self.driver.get(url)
            time.sleep(1)
            
            # Look for download button
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            download_link = soup.find('a', string=lambda x: x and 'download' in x.lower())
            
            if not download_link:
                download_link = soup.find('a', href=lambda x: x and 'download' in x)
            
            if download_link:
                download_url = download_link['href']
                if not download_url.startswith('http'):
                    download_url = urljoin(self.base_url, download_url)
            else:
                # Try current URL
                download_url = self.driver.current_url
            
            # Get filename from URL or page
            filename = None
            
            # Try to get filename from download link
            if '/' in download_url:
                filename = download_url.split('/')[-1].split('?')[0]
                filename = requests.utils.unquote(filename)
            
            # Try to get from page title
            if not filename or len(filename) < 3:
                title = soup.find('title')
                if title:
                    filename = title.get_text(strip=True)
                    filename = self._normalize_filename(filename)
            
            if not filename:
                filename = f"file_{len(self.visited_files)}"
            
            # Ensure filename has extension
            if not any(filename.endswith(ext) for ext in ['.pdf', '.pptx', '.docx', '.doc', '.zip', '.txt', '.html']):
                # Try to guess from content-type
                filename += '.unknown'
            
            file_path = save_dir / filename
            
            # Download using requests with cookies from selenium
            cookies = self.driver.get_cookies()
            session = requests.Session()
            for cookie in cookies:
                session.cookies.set(cookie['name'], cookie['value'])
            
            response = session.get(download_url, stream=True)
            response.raise_for_status()
            
            with open(file_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            
            # Extract text from file if possible
            self._extract_text_from_file(file_path)
            
            self.content_index['files'].append({
                'filename': filename,
                'path': str(file_path.relative_to(self.output_dir)),
                'url': url
            })
            
            return file_path
            
        except Exception as e:
            print(f"      ✗ Error downloading file from {url}: {e}")
            return None
    
    def _extract_text_from_file(self, file_path):
        """Extract text content from downloaded files."""
        try:
            file_ext = file_path.suffix.lower()
            text_path = file_path.with_suffix(file_path.suffix + '.txt')
            
            if file_ext == '.pdf':
                self._extract_pdf_text(file_path, text_path)
            elif file_ext == '.pptx':
                self._extract_pptx_text(file_path, text_path)
            elif file_ext in ['.docx', '.doc']:
                self._extract_docx_text(file_path, text_path)
            
        except Exception as e:
            print(f"      Warning: Could not extract text from {file_path.name}: {e}")
    
    def _extract_pdf_text(self, pdf_path, output_path):
        """Extract text from PDF file."""
        try:
            import pdfplumber
            
            with pdfplumber.open(pdf_path) as pdf:
                text = ""
                for page in pdf.pages:
                    text += page.extract_text() or ""
                    text += "\n\n--- Page Break ---\n\n"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
                
        except Exception as e:
            # Fallback to PyPDF2
            try:
                import PyPDF2
                with open(pdf_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() or ""
                        text += "\n\n--- Page Break ---\n\n"
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
            except:
                raise e
    
    def _extract_pptx_text(self, pptx_path, output_path):
        """Extract text from PowerPoint file."""
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n\n=== Slide {slide_num} ===\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    def _extract_docx_text(self, docx_path, output_path):
        """Extract text from Word document."""
        from docx import Document
        
        doc = Document(docx_path)
        text = ""
        
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
    
    def _scrape_page(self, url, save_dir, page_name):
        """Scrape a Canvas page."""
        try:
            if url in self.visited_urls:
                return None
            
            self.visited_urls.add(url)
            
            self.driver.get(url)
            time.sleep(1)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            page_name = self._normalize_filename(page_name)
            page_dir = save_dir / page_name
            page_dir.mkdir(exist_ok=True)
            
            # Save HTML
            with open(page_dir / "page.html", 'w', encoding='utf-8') as f:
                f.write(self.driver.page_source)
            
            # Extract text
            content = soup.get_text(separator='\n', strip=True)
            with open(page_dir / "page.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Look for embedded files/links
            links = soup.find_all('a', href=True)
            for link in links:
                href = link['href']
                if '/files/' in href:
                    if not href.startswith('http'):
                        href = urljoin(self.base_url, href)
                    self._download_file(href, page_dir)
            
            return page_dir
            
        except Exception as e:
            print(f"      ✗ Error scraping page: {e}")
            return None
    
    def _scrape_assignment(self, url, save_dir, assign_name):
        """Scrape an assignment page."""
        try:
            if url in self.visited_urls:
                return None
            
            self.visited_urls.add(url)
            
            self.driver.get(url)
            time.sleep(1)
            
            soup = BeautifulSoup(self.driver.page_source, 'html.parser')
            
            assign_name = self._normalize_filename(assign_name)
            assign_dir = save_dir / assign_name
            assign_dir.mkdir(exist_ok=True)
            
            # Save HTML
            with open(assign_dir / "assignment.html", 'w', encoding='utf-8') as f:
                f.write(self.driver.page_source)
            
            # Extract text
            content = soup.get_text(separator='\n', strip=True)
            with open(assign_dir / "assignment.txt", 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Look for attached files
            links = soup.find_all('a', href=True)
            for link in links:
                href = link['href']
                if '/files/' in href:
                    if not href.startswith('http'):
                        href = urljoin(self.base_url, href)
                    self._download_file(href, assign_dir)
            
            return assign_dir
            
        except Exception as e:
            print(f"      ✗ Error scraping assignment: {e}")
            return None
    
    def _save_content_index(self):
        """Save the content index JSON file."""
        index_path = self.output_dir / "content_index.json"
        with open(index_path, 'w', encoding='utf-8') as f:
            json.dump(self.content_index, f, indent=2)
        print(f"\n✓ Content index saved to: {index_path}")


def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("Usage: python canvas_course_scraper.py <course_url>")
        print("Example: python canvas_course_scraper.py https://canvas.cornell.edu/courses/80403")
        sys.exit(1)
    
    course_url = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "scraped_course"
    
    scraper = CanvasCourseScr aper(course_url, output_dir)
    scraper.scrape_all()


if __name__ == "__main__":
    main()
